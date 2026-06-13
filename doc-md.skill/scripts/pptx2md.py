#!/usr/bin/env python3
"""Convert PowerPoint file to Markdown."""

import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    print("Error: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


def pptx_to_markdown(input_path: str, output_path: str = None) -> str:
    """Convert pptx file to markdown."""
    input_file = Path(input_path)
    if not input_file.exists():
        raise FileNotFoundError(f"File not found: {input_path}")
    
    if output_path is None:
        output_path = input_file.with_suffix('.md')
    
    prs = Presentation(input_path)
    
    with open(output_path, 'w', encoding='utf-8') as f:
        for idx, slide in enumerate(prs.slides, 1):
            # Extract title (first shape with text, usually title)
            title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    title = shape.text.strip()
                    break
            
            f.write(f"# Slide {idx}: {title}\n\n")
            
            # Extract bullet points and text
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text and text != title:
                            level = paragraph.level
                            indent = "  " * level
                            f.write(f"{indent}- {text}\n")
            
            f.write("\n")
    
    return str(output_path)


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python pptx2md.py <input.pptx> [output.md]")
        sys.exit(1)
    
    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None
    
    try:
        result = pptx_to_markdown(input_file, output_file)
        print(f"Converted: {result}")
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)
